"""
知识库服务：创建知识库、添加文件并做 RAG 切分与向量化
"""
import asyncio
import io
import logging
from typing import List, Optional, Dict, Any, AsyncGenerator
from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, func, delete

from app.models.knowledge_base import KnowledgeBase, KnowledgeBaseFile
from app.models.file import File
from app.models.chunk import Chunk
from app.schemas.knowledge_base import (
    KnowledgeBaseCreate,
    KnowledgeBaseResponse,
    KnowledgeBaseListResponse,
    KnowledgeBaseFileItem,
    KnowledgeBaseFileListResponse,
    ChunkItem,
    ChunkListResponse,
)
from app.services.file_service import FileService
from app.services.embedding_service import get_embeddings, get_embedding, get_embedding_for_image
from app.services.vector_store import get_vector_client, chunk_id_to_vector_id
from app.services.ocr_service import extract_text_from_image
from app.core.config import settings


class KnowledgeBaseService:
    """知识库服务类"""
    
    def __init__(self, db: AsyncSession):
        self.db = db
    
    async def create_knowledge_base(
        self,
        kb_data: KnowledgeBaseCreate,
        user_id: int
    ) -> KnowledgeBase:
        """创建知识库"""
        kb = KnowledgeBase(
            user_id=user_id,
            name=kb_data.name,
            description=kb_data.description,
            chunk_size=getattr(kb_data, "chunk_size", None),
            chunk_overlap=getattr(kb_data, "chunk_overlap", None),
            chunk_max_expand_ratio=str(kb_data.chunk_max_expand_ratio) if getattr(kb_data, "chunk_max_expand_ratio", None) is not None else None,
        )
        self.db.add(kb)
        await self.db.commit()
        await self.db.refresh(kb)
        return kb
    
    async def get_knowledge_bases(
        self,
        user_id: int,
        page: int = 1,
        page_size: int = 20
    ) -> KnowledgeBaseListResponse:
        """获取知识库列表"""
        offset = (page - 1) * page_size
        
        count_result = await self.db.execute(
            select(func.count()).select_from(KnowledgeBase).where(KnowledgeBase.user_id == user_id)
        )
        total = count_result.scalar()
        
        result = await self.db.execute(
            select(KnowledgeBase)
            .where(KnowledgeBase.user_id == user_id)
            .order_by(KnowledgeBase.created_at.desc())
            .offset(offset)
            .limit(page_size)
        )
        kbs = result.scalars().all()
        
        return KnowledgeBaseListResponse(
            knowledge_bases=[KnowledgeBaseResponse.model_validate(kb) for kb in kbs],
            total=total,
            page=page,
            page_size=page_size
        )
    
    async def get_knowledge_base(self, kb_id: int, user_id: int) -> Optional[KnowledgeBase]:
        """获取知识库"""
        result = await self.db.execute(
            select(KnowledgeBase).where(
                KnowledgeBase.id == kb_id,
                KnowledgeBase.user_id == user_id
            )
        )
        return result.scalar_one_or_none()
    
    async def update_knowledge_base(
        self,
        kb_id: int,
        kb_data: KnowledgeBaseCreate,
        user_id: int
    ) -> Optional[KnowledgeBase]:
        """更新知识库"""
        kb = await self.get_knowledge_base(kb_id, user_id)
        if not kb:
            return None
        
        kb.name = kb_data.name
        kb.description = kb_data.description
        if hasattr(kb_data, "chunk_size"):
            kb.chunk_size = kb_data.chunk_size
        if hasattr(kb_data, "chunk_overlap"):
            kb.chunk_overlap = kb_data.chunk_overlap
        if hasattr(kb_data, "chunk_max_expand_ratio"):
            kb.chunk_max_expand_ratio = str(kb_data.chunk_max_expand_ratio) if kb_data.chunk_max_expand_ratio is not None else None
        await self.db.commit()
        await self.db.refresh(kb)
        return kb
    
    async def delete_knowledge_base(self, kb_id: int, user_id: int) -> None:
        """删除知识库，包括清理 Milvus 中的向量数据"""
        import logging
        from sqlalchemy import delete
        
        kb = await self.get_knowledge_base(kb_id, user_id)
        if not kb:
            raise ValueError("知识库不存在")
        
        # 1. 查询该知识库的所有 chunks，获取 vector_ids
        chunks_result = await self.db.execute(
            select(Chunk).where(Chunk.knowledge_base_id == kb_id)
        )
        chunks = list(chunks_result.scalars().all())
        
        # 2. 从 Milvus 中删除对应的向量
        if chunks:
            try:
                vector_store = get_vector_client()
                vector_ids_to_delete = []
                for chunk in chunks:
                    # 使用确定性算法计算 vector_id（与插入时一致）
                    vid = int(chunk_id_to_vector_id(chunk.id))
                    vector_ids_to_delete.append(vid)
                
                if vector_ids_to_delete:
                    try:
                        # 检查集合是否存在
                        if vector_store.client.has_collection(vector_store._collection):
                            # 从 Milvus 中删除向量
                            vector_store.client.delete(
                                collection_name=vector_store._collection,
                                ids=vector_ids_to_delete
                            )
                            logging.info(f"从 Milvus 中删除了 {len(vector_ids_to_delete)} 个向量")
                        else:
                            logging.warning(f"Milvus 集合 {vector_store._collection} 不存在，跳过向量删除")
                    except Exception as e:
                        logging.error(f"删除 Milvus 向量失败: {e}，继续删除数据库记录")
            except Exception as e:
                logging.error(f"清理 Milvus 向量时出错: {e}，继续删除数据库记录")
        
        # 3. 删除数据库中的 chunks（级联删除会自动处理，但显式删除更清晰）
        if chunks:
            await self.db.execute(
                delete(Chunk).where(Chunk.knowledge_base_id == kb_id)
            )
            await self.db.flush()
        
        # 4. 删除知识库文件关联（级联删除会自动处理，但显式删除更清晰）
        await self.db.execute(
            delete(KnowledgeBaseFile).where(KnowledgeBaseFile.knowledge_base_id == kb_id)
        )
        await self.db.flush()
        
        # 5. 删除知识库本身
        await self.db.delete(kb)
        await self.db.commit()
        
        logging.info(f"成功删除知识库 {kb_id} 及其所有相关数据（包括 {len(chunks)} 个 chunks 和对应的向量）")
    
    @staticmethod
    def _extract_text(content: bytes, file_type: str) -> str:
        """从文件内容提取纯文本（支持 txt、pdf、docx、pptx、xlsx）"""
        ft = (file_type or "").lower()
        if ft == "txt":
            return content.decode("utf-8", errors="ignore").strip()
        if ft == "pdf":
            text = ""
            try:
                from PyPDF2 import PdfReader
                reader = PdfReader(io.BytesIO(content))
                text = "\n".join(
                    (page.extract_text() or "").strip()
                    for page in reader.pages
                ).strip()
            except Exception as e:
                logging.warning(f"PyPDF2 提取 PDF 失败: {e}")
            if not text:
                try:
                    import pdfplumber
                    with pdfplumber.open(io.BytesIO(content)) as pdf:
                        parts = []
                        for page in pdf.pages:
                            t = page.extract_text()
                            if t and t.strip():
                                parts.append(t.strip())
                        text = "\n".join(parts).strip() if parts else ""
                    if text:
                        logging.info("PDF 文本由 pdfplumber 提取（PyPDF2 未提取到内容）")
                except Exception as e:
                    logging.warning(f"pdfplumber 提取 PDF 失败: {e}")
            table_text = KnowledgeBaseService._extract_pdf_tables_static(content)
            if table_text:
                text = (text + "\n\n" + table_text).strip()
            return text
        if ft == "docx":
            try:
                from docx import Document
                doc = Document(io.BytesIO(content))
                parts = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        parts.append(para.text.strip())
                for table in doc.tables:
                    for row in table.rows:
                        for cell in row.cells:
                            if cell.text.strip():
                                parts.append(cell.text.strip())
                return "\n".join(parts).strip() if parts else ""
            except Exception as e:
                logging.warning(f"docx 文本提取失败: {e}")
                return ""
        if ft == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(io.BytesIO(content))
                parts = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text and shape.text.strip():
                            parts.append(shape.text.strip())
                        # 表格
                        if shape.has_table:
                            for row in shape.table.rows:
                                for cell in row.cells:
                                    if cell.text and cell.text.strip():
                                        parts.append(cell.text.strip())
                return "\n".join(parts).strip() if parts else ""
            except Exception as e:
                logging.warning(f"pptx 文本提取失败: {e}")
                return ""
        if ft == "xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
                parts = []
                for name in wb.sheetnames:
                    sheet = wb[name]
                    sheet_parts = [f"表：{name}"]
                    for row in sheet.iter_rows(values_only=True):
                        row_str = "\t".join(str(c) if c is not None else "" for c in row).strip()
                        if row_str:
                            sheet_parts.append(row_str)
                    if len(sheet_parts) > 1:
                        parts.append("\n".join(sheet_parts))
                wb.close()
                return "\n\n".join(parts).strip() if parts else ""
            except Exception as e:
                logging.warning(f"xlsx 文本提取失败: {e}")
                return ""
        if ft in ("md", "markdown"):
            return content.decode("utf-8", errors="ignore").strip()
        if ft == "html":
            try:
                from bs4 import BeautifulSoup
                soup = BeautifulSoup(content.decode("utf-8", errors="ignore"), "html.parser")
                for tag in soup(["script", "style"]):
                    tag.decompose()
                return soup.get_text(separator="\n", strip=True) or ""
            except Exception as e:
                logging.warning(f"HTML 文本提取失败: {e}")
                return content.decode("utf-8", errors="ignore").strip()
        if ft == "zip":
            return KnowledgeBaseService._extract_zip_static(content)
        return ""

    @staticmethod
    def _extract_pdf_tables_static(content: bytes) -> str:
        """从 PDF 中提取表格，格式为「表：第N页表格」+ 行列文本，便于查表类问答。"""
        try:
            import pdfplumber
            parts = []
            with pdfplumber.open(io.BytesIO(content)) as pdf:
                for i, page in enumerate(pdf.pages):
                    tables = page.extract_tables()
                    for j, table in enumerate(tables or []):
                        if not table:
                            continue
                        title = f"表：第{i+1}页表格{j+1}"
                        rows = ["\t".join(str(cell or "").strip() for cell in row) for row in table]
                        if rows:
                            parts.append(title + "\n" + "\n".join(rows))
            return "\n\n".join(parts).strip() if parts else ""
        except Exception as e:
            logging.warning(f"PDF 表格提取失败: {e}")
            return ""

    @staticmethod
    def _extract_zip_static(content: bytes) -> str:
        """从 zip 中解压支持格式的文件，逐个提取文本后合并（带 [文件: 名] 前缀）。"""
        import zipfile
        supported = {"txt", "pdf", "md", "markdown", "html", "docx", "pptx", "xlsx"}
        parts = []
        try:
            with zipfile.ZipFile(io.BytesIO(content), "r") as z:
                for name in z.namelist():
                    if name.startswith("__MACOSX") or "/." in name:
                        continue
                    ext = name.split(".")[-1].lower() if "." in name else ""
                    if ext not in supported:
                        continue
                    try:
                        raw = z.read(name)
                        t = KnowledgeBaseService._extract_text(raw, ext if ext != "markdown" else "md")
                        if t and t.strip():
                            parts.append(f"[文件: {name}]\n{t}")
                    except Exception as e:
                        logging.warning(f"zip 内文件 {name} 提取失败: {e}")
            return "\n\n".join(parts).strip() if parts else ""
        except Exception as e:
            logging.warning(f"ZIP 解压/解析失败: {e}")
            return ""

    async def _extract_pdf_ocr(self, content: bytes) -> str:
        """扫描版 PDF：将每页渲染为图后走 OCR，再拼接文本。"""
        min_chars = getattr(settings, "PDF_OCR_MIN_CHARS", 80)
        dpi = getattr(settings, "PDF_OCR_DPI", 150)
        parts = []
        try:
            from pdf2image import convert_from_bytes
            images = convert_from_bytes(content, dpi=dpi)
            for img in images:
                buf = io.BytesIO()
                img.save(buf, "PNG")
                t = await extract_text_from_image(buf.getvalue(), "png")
                parts.append(t or "")
            return "\n\n".join(parts).strip()
        except ImportError:
            logging.warning("pdf2image 未安装，无法对 PDF 做 OCR（需安装 poppler）")
            return ""
        except Exception as e:
            logging.warning(f"PDF OCR 失败: {e}")
            return ""

    @staticmethod
    def _chunk_text(text: str, chunk_size: int = 500, overlap: int = 50, max_expand_ratio: float = 1.3) -> List[str]:
        """智能分块：按句子切分，保持语义完整性，不截断句子。
        
        Args:
            text: 待切分的文本
            chunk_size: 目标块大小（字符数）
            overlap: 重叠字符数
            max_expand_ratio: 最大扩展比例（允许超出 chunk_size 的最大倍数，默认 1.3）
        
        Returns:
            List[str]: 切分后的文本块列表
        """
        if not text or chunk_size <= 0:
            return []
        
        import re
        
        # 1. 按句子分割（使用中文和英文标点符号）
        # 匹配句号、问号、感叹号、换行符等作为句子结束符
        # 使用正则表达式分割，保留分隔符
        sentence_pattern = r'([。！？\n]+|[.!?\n]+)'
        parts = re.split(sentence_pattern, text)
        sentences = []
        current_sentence = ""
        for i, part in enumerate(parts):
            if re.match(sentence_pattern, part):
                # 这是分隔符，结束当前句子
                if current_sentence.strip():
                    sentences.append(current_sentence.strip() + part.strip())
                current_sentence = ""
            else:
                current_sentence += part
        # 添加最后一个句子（如果没有以分隔符结尾）
        if current_sentence.strip():
            sentences.append(current_sentence.strip())
        
        # 过滤空句子
        sentences = [s for s in sentences if s.strip()]
        
        if not sentences:
            # 如果没有找到句子分隔符，按段落分割
            paragraphs = text.split('\n\n')
            if len(paragraphs) == 1:
                # 如果也没有段落，按固定长度切分（兜底）
                chunks = []
                start = 0
                while start < len(text):
                    end = start + chunk_size
                    chunks.append(text[start:end])
                    start = end - overlap
                    if start >= len(text):
                        break
                return chunks
            sentences = [p.strip() for p in paragraphs if p.strip()]
        
        if not sentences:
            return []
        
        # 2. 智能合并句子，形成语义完整的块
        chunks = []
        current_chunk = []
        current_length = 0
        max_chunk_size = int(chunk_size * max_expand_ratio)  # 允许的最大块大小
        
        for i, sentence in enumerate(sentences):
            sentence_length = len(sentence)
            
            # 如果单个句子就超过最大块大小，需要进一步切分（但保持句子内不截断）
            if sentence_length > max_chunk_size:
                # 如果当前块有内容，先保存
                if current_chunk:
                    chunks.append(' '.join(current_chunk))
                    current_chunk = []
                    current_length = 0
                
                # 对超长句子，尝试按逗号、分号等进一步分割
                sub_sentences = re.split(r'[，；,;]+', sentence)
                sub_sentences = [s.strip() for s in sub_sentences if s.strip()]
                
                for sub_sentence in sub_sentences:
                    sub_length = len(sub_sentence)
                    if current_length + sub_length <= max_chunk_size:
                        current_chunk.append(sub_sentence)
                        current_length += sub_length + 1  # +1 是空格
                    else:
                        if current_chunk:
                            chunks.append(' '.join(current_chunk))
                        # 如果子句仍然太长，直接作为一个块
                        if sub_length > max_chunk_size:
                            chunks.append(sub_sentence)
                            current_chunk = []
                            current_length = 0
                        else:
                            current_chunk = [sub_sentence]
                            current_length = sub_length
                continue
            
            # 检查添加当前句子后是否会超出限制
            # 计算新长度：当前长度 + 句子长度 + 空格（如果有已有句子）
            separator_length = 1 if current_chunk else 0
            new_length = current_length + sentence_length + separator_length
            
            if new_length <= chunk_size:
                # 在目标大小内，直接添加
                current_chunk.append(sentence)
                current_length = new_length
            elif new_length <= max_chunk_size:
                # 超出目标大小但在允许范围内，为了保持语义完整性，仍然添加
                current_chunk.append(sentence)
                current_length = new_length
            else:
                # 超出允许范围，保存当前块，开始新块
                # 先提取重叠部分（在保存当前块之前）
                overlap_sentences = []
                overlap_length = 0
                if current_chunk and len(current_chunk) > 1:
                    # 从后往前取句子，直到达到重叠大小
                    for j in range(len(current_chunk) - 1, -1, -1):
                        sent = current_chunk[j]
                        sent_len = len(sent)
                        if overlap_length + sent_len <= overlap:
                            overlap_sentences.insert(0, sent)
                            overlap_length += sent_len + 1  # +1 是空格
                        else:
                            break
                
                # 保存当前块
                if current_chunk:
                    chunks.append(' '.join(current_chunk))
                
                # 开始新块，使用重叠句子
                current_chunk = overlap_sentences + [sentence]
                # 计算新块长度：重叠部分长度 + 当前句子长度 + 空格
                separator_count = len(overlap_sentences)  # 重叠句子之间的空格数
                current_length = overlap_length + sentence_length + separator_count
        
        # 添加最后一个块
        if current_chunk:
            chunks.append(' '.join(current_chunk))
        
        return chunks

    def _get_chunk_params(self, kb: Optional[KnowledgeBase], file_type: Optional[str] = None) -> tuple:
        """从知识库（及可选文件类型）解析分块参数，未设置则用全局 config。返回 (chunk_size, overlap, max_expand_ratio)。"""
        chunk_size = getattr(kb, "chunk_size", None) if kb else None
        chunk_overlap = getattr(kb, "chunk_overlap", None) if kb else None
        ratio_raw = getattr(kb, "chunk_max_expand_ratio", None) if kb else None
        if chunk_size is None:
            chunk_size = settings.CHUNK_SIZE
        if chunk_overlap is None:
            chunk_overlap = settings.CHUNK_OVERLAP
        try:
            max_expand_ratio = float(ratio_raw) if ratio_raw is not None else settings.CHUNK_MAX_EXPAND_RATIO
        except (TypeError, ValueError):
            max_expand_ratio = settings.CHUNK_MAX_EXPAND_RATIO
        return (chunk_size, chunk_overlap, max_expand_ratio)

    async def add_files(
        self, kb_id: int, file_ids: List[int], user_id: int
    ) -> tuple[Optional[KnowledgeBase], List[Dict[str, Any]]]:
        """添加文件到知识库并执行 RAG 切分与向量化。返回 (知识库, 被跳过的文件列表)。"""
        import logging
        from sqlalchemy.exc import SQLAlchemyError

        skipped: List[Dict[str, Any]] = []
        kb = await self.get_knowledge_base(kb_id, user_id)
        if not kb:
            return None, []

        file_service = FileService(self.db)
        vector_store = get_vector_client()
        
        # 在开始处理文件之前，先获取一个向量来确定实际维度
        actual_dim = None
        try:
            test_embedding = await get_embeddings(["test"])
            if test_embedding and len(test_embedding) > 0:
                actual_dim = len(test_embedding[0])
                logging.info(f"检测到向量维度: {actual_dim}")
        except Exception as e:
            logging.warning(f"无法预先获取向量维度: {e}，将使用配置的维度")
        
        try:
            vector_store.ensure_collection(actual_dim=actual_dim)
            logging.info(f"向量集合已确保存在，准备处理文件")
        except Exception as e:
            logging.error(f"创建向量集合失败: {e}")
            raise ValueError(f"无法创建向量集合，请检查 Zilliz 配置: {e}")

        added_kb_files = []
        added_chunks = []
        updated_files = []
        
        try:
            for file_id in file_ids:
                file_result = await self.db.execute(
                    select(File).where(File.id == file_id, File.user_id == user_id)
                )
                file = file_result.scalar_one_or_none()
                if not file:
                    continue

                existing_result = await self.db.execute(
                    select(KnowledgeBaseFile).where(
                        KnowledgeBaseFile.knowledge_base_id == kb_id,
                        KnowledgeBaseFile.file_id == file_id,
                    )
                )
                if existing_result.scalars().first():
                    continue

                kb_file = KnowledgeBaseFile(knowledge_base_id=kb_id, file_id=file_id)
                self.db.add(kb_file)
                await self.db.flush()
                added_kb_files.append(kb_file)

                content, content_error = await file_service.get_file_content(file_id, user_id)
                if not content:
                    logging.warning(f"文件 {file_id} 无法读取: {content_error}")
                    await self.db.delete(kb_file)
                    await self.db.flush()
                    skipped.append({
                        "file_id": file_id,
                        "original_filename": file.original_filename or file.filename,
                        "reason": content_error or "内容为空",
                    })
                    continue

                ft = (file.file_type or "").lower()
                if ft in ("jpeg", "jpg", "png"):
                    text = await extract_text_from_image(content, file.file_type)
                else:
                    text = self._extract_text(content, file.file_type)
                if ft == "pdf" and (not text or len(text.strip()) < getattr(settings, "PDF_OCR_MIN_CHARS", 80)):
                    ocr_text = await self._extract_pdf_ocr(content)
                    if ocr_text:
                        text = ocr_text
                        logging.info(f"文件 {file_id} 经 PDF OCR 补充文本")
                if not text or not text.strip():
                    logging.warning(f"文件 {file_id} 提取文本为空，跳过")
                    await self.db.delete(kb_file)
                    await self.db.flush()
                    skipped.append({
                        "file_id": file_id,
                        "original_filename": file.original_filename or file.filename,
                        "reason": "提取文本为空（可能为扫描版 PDF 或格式不支持）",
                    })
                    continue
                cs, co, ratio = self._get_chunk_params(kb, file.file_type)
                text_chunks = self._chunk_text(text, chunk_size=cs, overlap=co, max_expand_ratio=ratio)
                if not text_chunks:
                    logging.warning(f"文件 {file_id} 切分后无文本块，跳过")
                    await self.db.delete(kb_file)
                    await self.db.flush()
                    skipped.append({
                        "file_id": file_id,
                        "original_filename": file.original_filename or file.filename,
                        "reason": "切分后无文本块",
                    })
                    continue

                # 创建 Chunk 记录
                chunks = []
                for idx, chunk_text in enumerate(text_chunks):
                    chunk = Chunk(
                        file_id=file_id,
                        knowledge_base_id=kb_id,
                        content=chunk_text,
                        chunk_index=idx,
                    )
                    self.db.add(chunk)
                    chunks.append(chunk)
                await self.db.flush()
                added_chunks.extend(chunks)

                # 生成向量
                try:
                    embeddings = await get_embeddings([c.content for c in chunks])
                    if len(embeddings) != len(chunks):
                        raise ValueError(f"向量数量 {len(embeddings)} 与文本块数量 {len(chunks)} 不匹配")
                    
                    # 使用确定性 vector_id，与 vector_store 一致，供检索时反查
                    metadatas = []
                    for c, emb in zip(chunks, embeddings):
                        c.vector_id = chunk_id_to_vector_id(c.id)
                        meta = {
                            "chunk_id": c.id,
                            "content": c.content[:1000],
                            "file_id": c.file_id,
                            "knowledge_base_id": c.knowledge_base_id,
                            "chunk_index": c.chunk_index,
                            "embedding_source": "text",
                        }
                        metadatas.append(meta)
                    
                    ids_list = [str(c.id) for c in chunks]
                    await asyncio.get_event_loop().run_in_executor(
                        None,
                        lambda: vector_store.insert(ids=ids_list, vectors=embeddings, metadatas=metadatas),
                    )
                    logging.info(f"成功插入 {len(chunks)} 个向量到向量库（包含原文）")

                    # 图片文件：额外写入图像向量，支持图搜图、以文搜图统一空间
                    if ft in ("jpeg", "jpg", "png"):
                        try:
                            img_chunk = Chunk(
                                file_id=file_id,
                                knowledge_base_id=kb_id,
                                content=(text[:2000] if text else "[图片]"),
                                chunk_index=len(chunks),
                                chunk_metadata={"embedding_source": "image"},
                            )
                            self.db.add(img_chunk)
                            await self.db.flush()
                            added_chunks.append(img_chunk)
                            img_vec = await get_embedding_for_image(content, ft.replace("jpg", "jpeg"))
                            img_chunk.vector_id = chunk_id_to_vector_id(img_chunk.id)
                            img_meta = {
                                "chunk_id": img_chunk.id,
                                "content": (text[:500] if text else "[图片]"),
                                "file_id": file_id,
                                "knowledge_base_id": kb_id,
                                "chunk_index": img_chunk.chunk_index,
                                "embedding_source": "image",
                            }
                            await asyncio.get_event_loop().run_in_executor(
                                None,
                                lambda: vector_store.insert(
                                    ids=[str(img_chunk.id)],
                                    vectors=[img_vec],
                                    metadatas=[img_meta],
                                ),
                            )
                        except Exception as img_e:
                            logging.warning(f"文件 {file_id} 图像向量写入失败（已保留文本块）: {img_e}")
                    
                except Exception as e:
                    logging.error(f"文件 {file_id} 向量化失败: {e}")
                    raise ValueError(f"文件 {file_id} 向量化失败: {e}")

                old_chunk_count = file.chunk_count or 0
                file.chunk_count = old_chunk_count + len(chunks)
                if ft in ("jpeg", "jpg", "png") and any(
                    getattr(c, "chunk_metadata") and (c.chunk_metadata or {}).get("embedding_source") == "image"
                    for c in added_chunks if c.file_id == file_id
                ):
                    file.chunk_count += 1
                updated_files.append((file, old_chunk_count))

            # 更新知识库统计
            for file, old_count in updated_files:
                kb.chunk_count = (kb.chunk_count or 0) + (file.chunk_count - old_count)

            # 更新知识库 file_count
            count_result = await self.db.execute(
                select(func.count()).select_from(KnowledgeBaseFile).where(
                    KnowledgeBaseFile.knowledge_base_id == kb_id
                )
            )
            kb.file_count = count_result.scalar() or 0
            
            await self.db.commit()
            await self.db.refresh(kb)
            logging.info(f"成功处理 {len(added_kb_files)} 个文件，共 {len(added_chunks)} 个文本块；跳过 {len(skipped)} 个")
            return kb, skipped
            
        except Exception as e:
            # 发生错误，回滚所有操作
            logging.error(f"处理文件时发生错误: {e}，开始回滚")
            try:
                await self.db.rollback()
                logging.info("数据库事务已回滚")
            except Exception as rollback_error:
                logging.error(f"回滚失败: {rollback_error}")
            
            # 尝试清理已插入的向量（如果向量插入成功但后续失败）
            if added_chunks:
                try:
                    vector_ids_to_delete = [str(chunk_id_to_vector_id(c.id)) for c in added_chunks]
                    # 注意：这里需要根据实际的 vector_store 实现来删除向量
                    # 如果 vector_store 没有 delete 方法，可能需要手动调用 Milvus API
                    logging.warning(f"需要清理 {len(vector_ids_to_delete)} 个已插入的向量")
                except Exception as cleanup_error:
                    logging.error(f"清理向量失败: {cleanup_error}")
            
            raise ValueError(f"添加文件到知识库失败: {e}")

    async def add_files_stream(
        self, kb_id: int, file_ids: List[int], user_id: int
    ) -> AsyncGenerator[Dict[str, Any], None]:
        """添加文件到知识库（流式进度）。依次 yield file_start / file_done / file_skip，最后 yield done。"""
        import logging
        from sqlalchemy.exc import SQLAlchemyError

        skipped: List[Dict[str, Any]] = []
        kb = await self.get_knowledge_base(kb_id, user_id)
        if not kb:
            yield {"type": "error", "message": "知识库不存在"}
            return

        file_service = FileService(self.db)
        vector_store = get_vector_client()
        actual_dim = None
        try:
            test_embedding = await get_embeddings(["test"])
            if test_embedding and len(test_embedding) > 0:
                actual_dim = len(test_embedding[0])
        except Exception as e:
            logging.warning(f"无法预先获取向量维度: {e}")
        try:
            vector_store.ensure_collection(actual_dim=actual_dim)
        except Exception as e:
            yield {"type": "error", "message": str(e)}
            return

        added_kb_files = []
        added_chunks = []
        updated_files = []

        try:
            for file_id in file_ids:
                file_result = await self.db.execute(
                    select(File).where(File.id == file_id, File.user_id == user_id)
                )
                file = file_result.scalar_one_or_none()
                if not file:
                    yield {"type": "file_skip", "file_id": file_id, "filename": f"文件 {file_id}", "reason": "文件不存在或无权访问"}
                    continue
                filename = file.original_filename or file.filename or ""
                yield {"type": "file_start", "file_id": file_id, "filename": filename}

                existing_result = await self.db.execute(
                    select(KnowledgeBaseFile).where(
                        KnowledgeBaseFile.knowledge_base_id == kb_id,
                        KnowledgeBaseFile.file_id == file_id,
                    )
                )
                if existing_result.scalars().first():
                    skipped.append({"file_id": file_id, "original_filename": filename, "reason": "已在知识库中"})
                    yield {"type": "file_skip", "file_id": file_id, "filename": filename, "reason": "已在知识库中"}
                    continue

                kb_file = KnowledgeBaseFile(knowledge_base_id=kb_id, file_id=file_id)
                self.db.add(kb_file)
                await self.db.flush()
                added_kb_files.append(kb_file)

                content, content_error = await file_service.get_file_content(file_id, user_id)
                if not content:
                    await self.db.delete(kb_file)
                    await self.db.flush()
                    skipped.append({"file_id": file_id, "original_filename": filename, "reason": content_error or "内容为空"})
                    yield {"type": "file_skip", "file_id": file_id, "filename": filename, "reason": content_error or "内容为空"}
                    continue

                ft = (file.file_type or "").lower()
                if ft in ("jpeg", "jpg", "png"):
                    text = await extract_text_from_image(content, file.file_type)
                else:
                    text = self._extract_text(content, file.file_type)
                if ft == "pdf" and (not text or len(text.strip()) < getattr(settings, "PDF_OCR_MIN_CHARS", 80)):
                    ocr_text = await self._extract_pdf_ocr(content)
                    if ocr_text:
                        text = ocr_text
                if not text or not text.strip():
                    await self.db.delete(kb_file)
                    await self.db.flush()
                    skipped.append({"file_id": file_id, "original_filename": filename, "reason": "提取文本为空"})
                    yield {"type": "file_skip", "file_id": file_id, "filename": filename, "reason": "提取文本为空"}
                    continue
                cs, co, ratio = self._get_chunk_params(kb, file.file_type)
                text_chunks = self._chunk_text(text, chunk_size=cs, overlap=co, max_expand_ratio=ratio)
                if not text_chunks:
                    await self.db.delete(kb_file)
                    await self.db.flush()
                    skipped.append({"file_id": file_id, "original_filename": filename, "reason": "切分后无文本块"})
                    yield {"type": "file_skip", "file_id": file_id, "filename": filename, "reason": "切分后无文本块"}
                    continue

                chunks = []
                for idx, chunk_text in enumerate(text_chunks):
                    chunk = Chunk(
                        file_id=file_id,
                        knowledge_base_id=kb_id,
                        content=chunk_text,
                        chunk_index=idx,
                    )
                    self.db.add(chunk)
                    chunks.append(chunk)
                await self.db.flush()
                added_chunks.extend(chunks)

                try:
                    embeddings = await get_embeddings([c.content for c in chunks])
                    if len(embeddings) != len(chunks):
                        raise ValueError("向量数量与文本块数量不匹配")
                    metadatas = []
                    for c, emb in zip(chunks, embeddings):
                        c.vector_id = chunk_id_to_vector_id(c.id)
                        metadatas.append({
                            "chunk_id": c.id,
                            "content": c.content[:1000],
                            "file_id": c.file_id,
                            "knowledge_base_id": c.knowledge_base_id,
                            "chunk_index": c.chunk_index,
                            "embedding_source": "text",
                        })
                    ids_list = [str(c.id) for c in chunks]
                    await asyncio.get_event_loop().run_in_executor(
                        None,
                        lambda: vector_store.insert(ids=ids_list, vectors=embeddings, metadatas=metadatas),
                    )
                    extra_image_chunks = 0
                    if ft in ("jpeg", "jpg", "png"):
                        try:
                            img_chunk = Chunk(
                                file_id=file_id,
                                knowledge_base_id=kb_id,
                                content=(text[:2000] if text else "[图片]"),
                                chunk_index=len(chunks),
                                chunk_metadata={"embedding_source": "image"},
                            )
                            self.db.add(img_chunk)
                            await self.db.flush()
                            added_chunks.append(img_chunk)
                            img_vec = await get_embedding_for_image(content, ft.replace("jpg", "jpeg"))
                            img_chunk.vector_id = chunk_id_to_vector_id(img_chunk.id)
                            img_meta = {
                                "chunk_id": img_chunk.id,
                                "content": (text[:500] if text else "[图片]"),
                                "file_id": file_id,
                                "knowledge_base_id": kb_id,
                                "chunk_index": img_chunk.chunk_index,
                                "embedding_source": "image",
                            }
                            await asyncio.get_event_loop().run_in_executor(
                                None,
                                lambda: vector_store.insert(
                                    ids=[str(img_chunk.id)],
                                    vectors=[img_vec],
                                    metadatas=[img_meta],
                                ),
                            )
                            extra_image_chunks = 1
                        except Exception as img_e:
                            logging.warning(f"文件 {file_id} 图像向量写入失败: {img_e}")
                except Exception as e:
                    logging.error(f"文件 {file_id} 向量化失败: {e}")
                    await self.db.delete(kb_file)
                    await self.db.flush()
                    skipped.append({"file_id": file_id, "original_filename": filename, "reason": f"向量化失败: {e}"})
                    yield {"type": "file_skip", "file_id": file_id, "filename": filename, "reason": f"向量化失败: {str(e)}"}
                    continue

                old_chunk_count = file.chunk_count or 0
                file.chunk_count = old_chunk_count + len(chunks) + extra_image_chunks
                updated_files.append((file, old_chunk_count))
                yield {"type": "file_done", "file_id": file_id, "filename": filename, "chunk_count": len(chunks) + extra_image_chunks}

            for file, old_count in updated_files:
                kb.chunk_count = (kb.chunk_count or 0) + (file.chunk_count - old_count)
            count_result = await self.db.execute(
                select(func.count()).select_from(KnowledgeBaseFile).where(
                    KnowledgeBaseFile.knowledge_base_id == kb_id
                )
            )
            kb.file_count = count_result.scalar() or 0
            await self.db.commit()
            await self.db.refresh(kb)
            yield {
                "type": "done",
                "knowledge_base": KnowledgeBaseResponse.model_validate(kb).model_dump(mode="json"),
                "skipped": skipped,
            }
        except Exception as e:
            logging.exception("add_files_stream 失败")
            try:
                await self.db.rollback()
            except Exception:
                pass
            yield {"type": "error", "message": str(e)}

    async def get_files_in_knowledge_base(
        self,
        kb_id: int,
        user_id: int,
        page: int = 1,
        page_size: int = 20,
    ) -> KnowledgeBaseFileListResponse:
        """查询知识库内的文件列表（含该文件在本库中的分块数）"""
        kb = await self.get_knowledge_base(kb_id, user_id)
        if not kb:
            raise ValueError("知识库不存在")
        offset = (page - 1) * page_size
        # 总数：知识库内且属于当前用户的文件数
        total_result = await self.db.execute(
            select(func.count())
            .select_from(KnowledgeBaseFile)
            .join(File, KnowledgeBaseFile.file_id == File.id)
            .where(
                KnowledgeBaseFile.knowledge_base_id == kb_id,
                File.user_id == user_id,
            )
        )
        total = total_result.scalar() or 0
        # 列表：KnowledgeBaseFile join File
        result = await self.db.execute(
            select(KnowledgeBaseFile, File)
            .join(File, KnowledgeBaseFile.file_id == File.id)
            .where(
                KnowledgeBaseFile.knowledge_base_id == kb_id,
                File.user_id == user_id,
            )
            .order_by(KnowledgeBaseFile.created_at.desc())
            .offset(offset)
            .limit(page_size)
        )
        rows = result.all()
        items = []
        for kb_file, file in rows:
            chunk_count_result = await self.db.execute(
                select(func.count()).select_from(Chunk).where(
                    Chunk.knowledge_base_id == kb_id,
                    Chunk.file_id == file.id,
                )
            )
            chunk_count_in_kb = chunk_count_result.scalar() or 0
            items.append(
                KnowledgeBaseFileItem(
                    file_id=file.id,
                    original_filename=file.original_filename or file.filename,
                    file_type=file.file_type,
                    file_size=file.file_size,
                    chunk_count_in_kb=chunk_count_in_kb,
                    added_at=kb_file.created_at,
                )
            )
        return KnowledgeBaseFileListResponse(files=items, total=total, page=page, page_size=page_size)

    async def get_chunks_for_file_in_kb(
        self, kb_id: int, file_id: int, user_id: int
    ) -> ChunkListResponse:
        """查询某文件在知识库中的分块列表（按 chunk_index 排序）"""
        kb = await self.get_knowledge_base(kb_id, user_id)
        if not kb:
            raise ValueError("知识库不存在")
        kb_file_result = await self.db.execute(
            select(KnowledgeBaseFile).where(
                KnowledgeBaseFile.knowledge_base_id == kb_id,
                KnowledgeBaseFile.file_id == file_id,
            )
        )
        if not kb_file_result.scalar_one_or_none():
            raise ValueError("该文件不在本知识库中")
        file_result = await self.db.execute(select(File).where(File.id == file_id, File.user_id == user_id))
        if not file_result.scalar_one_or_none():
            raise ValueError("文件不存在或无权操作")
        result = await self.db.execute(
            select(Chunk)
            .where(Chunk.knowledge_base_id == kb_id, Chunk.file_id == file_id)
            .order_by(Chunk.chunk_index)
        )
        chunks = result.scalars().all()
        return ChunkListResponse(
            chunks=[ChunkItem(id=c.id, chunk_index=c.chunk_index, content=c.content or "") for c in chunks]
        )

    async def remove_file_from_knowledge_base(self, kb_id: int, file_id: int, user_id: int) -> None:
        """从知识库中移除文件：删除该文件在本库中的分块与向量，更新统计"""
        kb = await self.get_knowledge_base(kb_id, user_id)
        if not kb:
            raise ValueError("知识库不存在")
        kb_file_result = await self.db.execute(
            select(KnowledgeBaseFile).where(
                KnowledgeBaseFile.knowledge_base_id == kb_id,
                KnowledgeBaseFile.file_id == file_id,
            )
        )
        kb_file = kb_file_result.scalar_one_or_none()
        if not kb_file:
            raise ValueError("该文件不在本知识库中")
        file_result = await self.db.execute(select(File).where(File.id == file_id, File.user_id == user_id))
        file = file_result.scalar_one_or_none()
        if not file:
            raise ValueError("文件不存在或无权操作")
        chunks_result = await self.db.execute(
            select(Chunk).where(Chunk.knowledge_base_id == kb_id, Chunk.file_id == file_id)
        )
        chunks = list(chunks_result.scalars().all())
        vector_store = get_vector_client()
        if chunks:
            try:
                if vector_store.client.has_collection(vector_store._collection):
                    vector_ids = [int(chunk_id_to_vector_id(c.id)) for c in chunks]
                    vector_store.client.delete(collection_name=vector_store._collection, ids=vector_ids)
                    logging.info(f"从向量库删除了 {len(vector_ids)} 个向量")
            except Exception as e:
                logging.warning(f"删除向量失败: {e}，继续删除数据库记录")
        await self.db.execute(delete(Chunk).where(Chunk.knowledge_base_id == kb_id, Chunk.file_id == file_id))
        await self.db.flush()
        await self.db.delete(kb_file)
        await self.db.flush()
        chunk_delta = len(chunks)
        file.chunk_count = max(0, (file.chunk_count or 0) - chunk_delta)
        kb.file_count = max(0, (kb.file_count or 0) - 1)
        kb.chunk_count = max(0, (kb.chunk_count or 0) - chunk_delta)
        await self.db.commit()
        logging.info(f"已从知识库 {kb_id} 移除文件 {file_id}，删除 {chunk_delta} 个分块")

    async def reindex_file_in_knowledge_base(self, kb_id: int, file_id: int, user_id: int) -> Optional[KnowledgeBase]:
        """重新索引：先移除该文件在本库中的分块与向量，再重新切分与向量化"""
        await self.remove_file_from_knowledge_base(kb_id, file_id, user_id)
        kb, _ = await self.add_files(kb_id, [file_id], user_id)
        return kb

    async def search_images_by_text(
        self,
        query: str,
        user_id: int,
        knowledge_base_id: Optional[int] = None,
        top_k: int = 20,
    ) -> List[Dict[str, Any]]:
        """以文搜图：根据文本在知识库中检索匹配的图片文件。
        
        使用查询文本的向量在向量库中检索，再过滤出 file_type 为 jpeg/jpg/png 的文件，
        按相似度排序后去重（同一文件只返回一次），返回文件信息及片段。
        """
        if not (query and query.strip()):
            return []
        try:
            query_vec = await get_embedding(query.strip())
            vs = get_vector_client()
            filter_expr = f"knowledge_base_id == {knowledge_base_id}" if knowledge_base_id else None
            hits = vs.search(query_vector=query_vec, top_k=min(80, top_k * 4), filter_expr=filter_expr) or []
        except Exception as e:
            logging.warning(f"以文搜图向量检索失败: {e}")
            return []

        vector_ids = []
        for h in hits if isinstance(hits, list) else []:
            if not isinstance(h, dict):
                continue
            vid = h.get("id") or (h.get("entity") or {}).get("id") if isinstance(h.get("entity"), dict) else None
            if vid is not None:
                vector_ids.append(str(vid))
        if not vector_ids:
            return []

        # Chunk join File，只保留图片类型且属于当前用户的文件
        stmt = (
            select(Chunk, File)
            .join(File, Chunk.file_id == File.id)
            .where(
                Chunk.vector_id.in_(vector_ids),
                File.file_type.in_(("jpeg", "jpg", "png")),
                File.user_id == user_id,
            )
        )
        if knowledge_base_id is not None:
            stmt = stmt.where(Chunk.knowledge_base_id == knowledge_base_id)
        result = await self.db.execute(stmt)
        rows = result.all()
        # 按向量检索顺序排序，同一 file_id 只保留第一次出现（最佳匹配）
        seen_file_ids = set()
        ordered_files: List[Dict[str, Any]] = []
        vid_order = {vid: i for i, vid in enumerate(vector_ids)}
        for chunk, file in rows:
            if file.id in seen_file_ids:
                continue
            seen_file_ids.add(file.id)
            rank = vid_order.get(chunk.vector_id, 9999)
            ordered_files.append({
                "rank": rank,
                "file_id": file.id,
                "original_filename": file.original_filename or file.filename,
                "file_type": file.file_type,
                "snippet": (chunk.content or "")[:200],
            })
        ordered_files.sort(key=lambda x: x["rank"])
        return ordered_files[:top_k]

    async def search_unified(
        self,
        query: Optional[str] = None,
        image_bytes: Optional[bytes] = None,
        user_id: int = None,
        knowledge_base_id: Optional[int] = None,
        top_k: int = 30,
    ) -> List[Dict[str, Any]]:
        """多模态检索统一：以文或以图一次查询，同时返回文档与图片（同一向量空间）。
        若提供 image_bytes 则用图向量，否则用 query 文本向量。
        """
        if image_bytes:
            query_vec = await get_embedding_for_image(
                image_bytes,
                image_format="jpeg",
            )
        elif query and query.strip():
            query_vec = await get_embedding(query.strip())
        else:
            return []
        vs = get_vector_client()
        filter_expr = f"knowledge_base_id == {knowledge_base_id}" if knowledge_base_id else None
        hits = vs.search(query_vector=query_vec, top_k=min(80, top_k * 2), filter_expr=filter_expr) or []
        vector_ids = []
        id_to_score = {}
        for rank, h in enumerate(hits if isinstance(hits, list) else []):
            if not isinstance(h, dict):
                continue
            vid = h.get("id") or (h.get("entity") or h.get("payload") or {}).get("id")
            if vid is None:
                continue
            vid_str = str(vid)
            vector_ids.append(vid_str)
            dist = h.get("distance")
            if dist is not None:
                id_to_score[vid_str] = float(dist)
            else:
                id_to_score[vid_str] = 1.0 - (rank / max(len(hits), 1))
        if not vector_ids:
            return []
        stmt = (
            select(Chunk, File)
            .join(File, Chunk.file_id == File.id)
            .where(
                Chunk.vector_id.in_(vector_ids),
                File.user_id == user_id,
            )
        )
        if knowledge_base_id is not None:
            stmt = stmt.where(Chunk.knowledge_base_id == knowledge_base_id)
        result = await self.db.execute(stmt)
        rows = result.all()
        vid_to_rank = {vid: i for i, vid in enumerate(vector_ids)}
        items = []
        for chunk, file in rows:
            rank = vid_to_rank.get(str(chunk.vector_id), 9999)
            score = id_to_score.get(str(chunk.vector_id), 1.0 - rank / 100)
            items.append({
                "chunk_id": chunk.id,
                "file_id": file.id,
                "original_filename": file.original_filename or file.filename,
                "file_type": file.file_type or "",
                "snippet": (chunk.content or "")[:300],
                "score": score,
                "is_image": (file.file_type or "").lower() in ("jpeg", "jpg", "png"),
                "_rank": rank,
            })
        items.sort(key=lambda x: x["_rank"])
        for x in items:
            del x["_rank"]
        return items[:top_k]

    async def search_images_by_image(
        self,
        image_bytes: bytes,
        user_id: int,
        knowledge_base_id: Optional[int] = None,
        top_k: int = 20,
    ) -> List[Dict[str, Any]]:
        """图搜图：上传一张图，用多模态模型得到向量，在知识库中检索相似图片。
        返回所有命中向量且属于图片文件（file_type 为 jpeg/jpg/png）的记录；
        优先保留 embedding_source=image 的块（同一文件只取最佳排名的一条），
        这样既有图像向量的新图能高排，仅有 OCR 文本向量的旧图也有机会被搜到。
        """
        if not image_bytes or len(image_bytes) == 0:
            return []
        try:
            query_vec = await get_embedding_for_image(image_bytes, "jpeg")
        except Exception as e:
            logging.warning(f"图搜图获取图片向量失败: {e}")
            return []
        vs = get_vector_client()
        filter_expr = f"knowledge_base_id == {knowledge_base_id}" if knowledge_base_id else None
        hits = vs.search(query_vector=query_vec, top_k=min(100, top_k * 4), filter_expr=filter_expr) or []
        vector_ids = []
        id_to_score = {}
        for rank, h in enumerate(hits if isinstance(hits, list) else []):
            if not isinstance(h, dict):
                continue
            vid = h.get("id") or (h.get("entity") or h.get("payload") or {}).get("id")
            if vid is None:
                continue
            vid_str = str(vid)
            vector_ids.append(vid_str)
            dist = h.get("distance")
            id_to_score[vid_str] = float(dist) if dist is not None else (1.0 - rank / 100)
        if not vector_ids:
            return []
        stmt = (
            select(Chunk, File)
            .join(File, Chunk.file_id == File.id)
            .where(
                Chunk.vector_id.in_(vector_ids),
                File.user_id == user_id,
                File.file_type.in_(("jpeg", "jpg", "png")),
            )
        )
        if knowledge_base_id is not None:
            stmt = stmt.where(Chunk.knowledge_base_id == knowledge_base_id)
        result = await self.db.execute(stmt)
        rows = result.all()
        vid_order = {vid: i for i, vid in enumerate(vector_ids)}
        best_for_file: Dict[int, Dict[str, Any]] = {}
        for chunk, file in rows:
            rank = vid_order.get(str(chunk.vector_id), 9999)
            score = id_to_score.get(str(chunk.vector_id), 0.0)
            is_image_chunk = (chunk.chunk_metadata or {}).get("embedding_source") == "image"
            cand = {
                "file_id": file.id,
                "original_filename": file.original_filename or file.filename,
                "file_type": file.file_type,
                "snippet": (chunk.content or "")[:200],
                "score": score,
                "rank": rank,
                "_is_image_chunk": is_image_chunk,
                "_rank": rank,
            }
            existing = best_for_file.get(file.id)
            if existing is None:
                best_for_file[file.id] = cand
            else:
                replace = is_image_chunk and not existing.get("_is_image_chunk")
                if not replace and is_image_chunk == existing.get("_is_image_chunk") and rank < existing.get("_rank", 9999):
                    replace = True
                if replace:
                    best_for_file[file.id] = cand
        ordered = sorted(
            best_for_file.values(),
            key=lambda x: (not x.get("_is_image_chunk"), x.get("_rank", 9999)),
        )
        return [
            {"file_id": x["file_id"], "original_filename": x["original_filename"], "file_type": x["file_type"], "snippet": x.get("snippet"), "score": x["score"]}
            for x in ordered[:top_k]
        ]
